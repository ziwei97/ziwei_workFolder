import numpy as np
import sys
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum import text
from pptx.enum.text import PP_ALIGN
import pandas as pd

#test"
class AutoSlideGen:
    def __init__(self, imgLoc, maskKeyWord, saveTo):
        self.imgLoc = imgLoc
        self.maskKeyWord = maskKeyWord
        self.saveTo = saveTo
        self.total_num = 0

    def generate_slide(self):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]

        for dirs, subdir, files in os.walk(self.imgLoc):
            if not subdir:
                mask_imgs = [ele for ele in files if 'Mask_' in ele]  # 'Mask_Truth' in ele and 'HoleDeleted' in ele]

                if len(mask_imgs) > 0:

                    mask_img = mask_imgs[0]
                    truth_img = \
                    [ele for ele in files if ('.png' in ele and 'Mask' not in ele) and ele not in [mask_img]][0]
                    # pseudo_imgs = [ele for ele in files if ele not in [mask_img,truth_img]]
                    # pseudo_img = [ele for ele in files if 'PseudoColor' in ele and '.xcf' not in ele and '-raw' not in ele][0]
                    pseudo_img = [ele for ele in files if
                                  'PseudoColor' in ele and '.xcf' not in ele and '-raw' not in ele]
                    print(pseudo_img[0])

                    pseudo_img = os.path.join(dirs, pseudo_img[0])
                    truth_img = os.path.join(dirs, truth_img)
                    mask_img = os.path.join(dirs, mask_img)
                    # print(pseudo_img)
                    # print(mask_img)
                    # print(truth_img)
                    # create slide
                    slide = prs.slides.add_slide(slide_layout)

                    # add title
                    shapes = slide.shapes
                    shapes.title.text = os.path.basename(dirs)

                    # add images
                    left = Inches(0.1)
                    top = Inches(2)
                    height = Inches(3.2)
                    pseudo_pic = slide.shapes.add_picture(pseudo_img, left, top, height)

                    left = Inches(3.4)
                    truth_pic = slide.shapes.add_picture(truth_img, left, top, height)

                    left = Inches(6.7)
                    mask_pic = slide.shapes.add_picture(mask_img, left, top, height)

                    # add image title
                    left = Inches(0.1)
                    top = Inches(5.3)
                    width = Inches(3.2)
                    height = Inches(0.5)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = 'PseudoColor Image'
                    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                    left = Inches(3.4)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = 'Truth Image'
                    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                    left = Inches(6.7)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = 'Algorithm Mask'
                    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                    self.total_num += 1
                else:
                    print(self.imgLoc)

        prs.save(self.saveTo)


src_path = '/Users/ziweishi/Documents/new1/'
maskKeyWord = 'Mask_'
centers = os.listdir(src_path)
if '.DS_Store' in centers:
    centers.remove('.DS_Store')
# subjects = os.listdir(src_path)

dst_path = '/Users/ziweishi/Documents/review'
if not os.path.exists(dst_path):
    os.mkdir(dst_path)

mask_no = 0
for center in centers:
    sub_dir = os.path.join(src_path, center)
    subjects = os.listdir(sub_dir)
    for subject in subjects:
        imgLoc = os.path.join(sub_dir, subject)
        # imgLoc = os.path.join(src_path, subject)
        ppt_name = center + '_' + subject + '_Review.pptx'
        # ppt_name = subject + '_Review.pptx'
        saveTo = os.path.join(dst_path, ppt_name)
        asg = AutoSlideGen(imgLoc, maskKeyWord, saveTo)
        asg.generate_slide()
        # print('subject: {}, No: {}'.format(subject, asg.total_num))
        mask_no += asg.total_num
        print('Done...{}'.format(ppt_name))

print('Mask No：', mask_no)